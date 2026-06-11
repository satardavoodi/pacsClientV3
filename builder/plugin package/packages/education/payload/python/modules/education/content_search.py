"""content_search — full-content search across educational resources.

Qt-free, worker-thread safe (called by the Secretary background agent,
also usable synchronously). Searches:

* Courses DB        — name / description / author / tags
* Slides DB         — slide titles / notes / text content blocks
* Case of the Day   — reuses ``case_of_day_database.search_cases``
* Course asset FILES under ``EDUCATION_STORAGE_PATH`` — PDF / e-book /
  PPTX / DOCX / TXT, via OPTIONAL extractors (pypdf, python-pptx,
  python-docx). A missing extractor is reported in
  ``skipped_extractors`` — never an exception.
* Image-only PDFs   — OCR fallback when the pluggable OCR engine is
  available (see ``secretary.background.verification.ocr_available``).
* Consultation notes — ONLY when ``include_consultations=True``
  (permission stays with the caller/user).

File extraction results are cached in ``content_index.json`` keyed by
(path, mtime, size) so repeated searches don't re-parse unchanged files.
A per-file size cap and a cooperative ``cancelled`` callback keep the
scan bounded.
"""
from __future__ import annotations

import json
import logging
import re
import time
from pathlib import Path
from typing import Callable, Optional

logger = logging.getLogger(__name__)

MAX_FILE_MB = 50
MAX_CHARS_PER_FILE = 400_000
SNIPPET_RADIUS = 80

_TEXT_EXTS = {".txt", ".md", ".html", ".htm", ".json"}
_PDF_EXTS = {".pdf"}
_PPTX_EXTS = {".pptx"}
_DOCX_EXTS = {".docx"}


# ── storage roots ────────────────────────────────────────────────────────

def _education_root() -> Optional[Path]:
    try:
        from PacsClient.utils.config import EDUCATION_STORAGE_PATH
        p = Path(EDUCATION_STORAGE_PATH)
        return p if p.exists() else None
    except Exception:
        return None


def _index_path() -> Optional[Path]:
    root = _education_root()
    return (root / "content_index.json") if root else None


# ── extraction cache ─────────────────────────────────────────────────────

def _load_index() -> dict:
    p = _index_path()
    if p is None or not p.exists():
        return {}
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return {}


def _save_index(index: dict) -> None:
    p = _index_path()
    if p is None:
        return
    try:
        tmp = p.with_suffix(".json.tmp")
        tmp.write_text(json.dumps(index, ensure_ascii=False),
                       encoding="utf-8")
        tmp.replace(p)
    except Exception:
        logger.debug("content_search: index save failed", exc_info=True)


def _cache_key(path: Path) -> str:
    st = path.stat()
    return f"{path}|{int(st.st_mtime)}|{st.st_size}"


# ── extractors (all optional) ────────────────────────────────────────────

def _extract_pdf(path: Path) -> Optional[str]:
    try:
        from pypdf import PdfReader
    except ImportError:
        return None
    try:
        reader = PdfReader(str(path))
        parts: list[str] = []
        total = 0
        for page in reader.pages[:300]:
            t = page.extract_text() or ""
            parts.append(t)
            total += len(t)
            if total > MAX_CHARS_PER_FILE:
                break
        text = "\n".join(parts)
        if text.strip():
            return text
        # Image-only PDF → OCR fallback when available.
        return _ocr_pdf(path) or ""
    except Exception:
        logger.debug("content_search: pdf extract failed %s", path,
                     exc_info=True)
        return ""


def _ocr_pdf(path: Path) -> str:
    """OCR the first pages of an image-only PDF (pluggable engine)."""
    try:
        from modules.EchoMind.secretary.background import verification as V
        if not V.ocr_available():
            return ""
        from pypdf import PdfReader
        import io
        from PIL import Image
        reader = PdfReader(str(path))
        chunks: list[str] = []
        for page in reader.pages[:10]:
            for img in getattr(page, "images", [])[:4]:
                try:
                    with Image.open(io.BytesIO(img.data)) as im:
                        import pytesseract
                        pytesseract.pytesseract.tesseract_cmd = (
                            V._tesseract_binary())
                        chunks.append(pytesseract.image_to_string(im) or "")
                except Exception:
                    continue
        return "\n".join(chunks)
    except Exception:
        return ""


def _extract_pptx(path: Path) -> Optional[str]:
    try:
        from pptx import Presentation
    except ImportError:
        return None
    try:
        prs = Presentation(str(path))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    parts.append(shape.text_frame.text or "")
            if getattr(slide, "has_notes_slide", False):
                try:
                    parts.append(slide.notes_slide.notes_text_frame.text or "")
                except Exception:
                    pass
        return "\n".join(parts)
    except Exception:
        logger.debug("content_search: pptx extract failed %s", path,
                     exc_info=True)
        return ""


def _extract_docx(path: Path) -> Optional[str]:
    try:
        import docx  # python-docx
    except ImportError:
        return None
    try:
        d = docx.Document(str(path))
        return "\n".join(p.text for p in d.paragraphs)
    except Exception:
        return ""


def _extract_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="replace")[
            :MAX_CHARS_PER_FILE]
    except Exception:
        return ""


def extract_file_text(path: Path) -> tuple[Optional[str], str]:
    """Return (text|None, extractor_name). None text = extractor missing."""
    suffix = path.suffix.lower()
    if suffix in _TEXT_EXTS:
        return _extract_text_file(path), "text"
    if suffix in _PDF_EXTS:
        return _extract_pdf(path), "pdf"
    if suffix in _PPTX_EXTS:
        return _extract_pptx(path), "pptx"
    if suffix in _DOCX_EXTS:
        return _extract_docx(path), "docx"
    return "", "unsupported"


# ── matching helpers ─────────────────────────────────────────────────────

def _terms(query: str) -> list[str]:
    return [t for t in re.split(r"\W+", (query or "").lower()) if len(t) > 2] \
        or [t for t in re.split(r"\W+", (query or "").lower()) if t]


def _match(text: str, terms: list[str]) -> tuple[float, int]:
    """(ratio of terms present, position of first hit)."""
    if not text or not terms:
        return 0.0, -1
    low = text.lower()
    hits = [low.find(t) for t in terms]
    found = [h for h in hits if h >= 0]
    if not found:
        return 0.0, -1
    return len(found) / len(terms), min(found)


def _snippet(text: str, pos: int) -> str:
    if pos < 0:
        return text[:SNIPPET_RADIUS * 2].strip()
    start = max(0, pos - SNIPPET_RADIUS)
    end = min(len(text), pos + SNIPPET_RADIUS)
    prefix = "…" if start > 0 else ""
    suffix = "…" if end < len(text) else ""
    return prefix + re.sub(r"\s+", " ", text[start:end]).strip() + suffix


def _result(source: str, title: str, location: str, score: float,
            snippet: str, extra: Optional[dict] = None) -> dict:
    r = {"source": source, "title": title, "location": location,
         "score": round(score, 3), "snippet": snippet}
    if extra:
        r.update(extra)
    return r


# ── main entry ───────────────────────────────────────────────────────────

def search_education_content(
    query: str,
    include_consultations: bool = False,
    cancelled: Optional[Callable[[], bool]] = None,
    min_ratio: float = 0.5,
) -> dict:
    """Search every educational source; returns a structured report."""
    t0 = time.time()
    terms = _terms(query)
    cancelled = cancelled or (lambda: False)
    results: list[dict] = []
    skipped: list[str] = []
    sources = 0

    # 1) Courses DB (name/description/author/tags) ----------------------
    try:
        from modules.education.course_database import get_all_courses
        sources += 1
        for course in get_all_courses():
            blob = " ".join(str(course.get(k) or "") for k in
                            ("course_name", "name", "description", "author",
                             "tags", "modality", "resource_type"))
            ratio, pos = _match(blob, terms)
            if ratio >= min_ratio:
                results.append(_result(
                    "course",
                    str(course.get("course_name") or course.get("name") or "Course"),
                    f"course_pk={course.get('course_pk') or course.get('pk')}",
                    ratio, _snippet(blob, pos)))
    except Exception:
        logger.exception("content_search: course DB search failed")
        skipped.append("courses DB (error)")

    # 2) Slides DB (titles/notes/text blocks) ---------------------------
    try:
        from modules.education.course_database import (
            get_all_courses, get_slides_for_course, get_content_for_slide)
        sources += 1
        for course in get_all_courses():
            if cancelled():
                break
            cpk = course.get("course_pk") or course.get("pk")
            if cpk is None:
                continue
            for slide in get_slides_for_course(int(cpk)):
                parts = [str(slide.get("slide_title") or ""),
                         str(slide.get("slide_notes") or "")]
                try:
                    for block in get_content_for_slide(int(slide["slide_pk"])):
                        if str(block.get("content_type")) == "text":
                            data = block.get("content_data")
                            if isinstance(data, str):
                                try:
                                    data = json.loads(data)
                                except Exception:
                                    data = {"text": data}
                            if isinstance(data, dict):
                                parts.append(str(data.get("text") or
                                                 data.get("content") or ""))
                except Exception:
                    pass
                blob = " ".join(parts)
                ratio, pos = _match(blob, terms)
                if ratio >= min_ratio:
                    results.append(_result(
                        "slide",
                        f"{course.get('course_name') or 'Course'} — "
                        f"slide {slide.get('slide_order')}: "
                        f"{slide.get('slide_title') or ''}".strip(),
                        f"course_pk={cpk};slide_pk={slide.get('slide_pk')}",
                        ratio, _snippet(blob, pos)))
    except Exception:
        logger.exception("content_search: slide DB search failed")
        skipped.append("slides DB (error)")

    # 3) Case of the Day -------------------------------------------------
    try:
        from modules.education.case_of_day_database import search_cases
        sources += 1
        for entry in search_cases(query=query):
            blob = " ".join(str(getattr(entry, f, "") or "") for f in
                            ("diagnosis", "description",
                             "differential_diagnosis", "study_description"))
            ratio, pos = _match(blob, terms)
            results.append(_result(
                "case_of_day",
                getattr(entry, "diagnosis", "") or "Case of the Day",
                f"case_pk={getattr(entry, 'case_pk', '')}",
                max(ratio, 0.5), _snippet(blob, pos)))
    except Exception:
        logger.exception("content_search: case-of-day search failed")
        skipped.append("case of the day (error)")

    # 4) Files (PDF / e-book / PPTX / DOCX / TXT) ------------------------
    root = _education_root()
    if root is not None:
        sources += 1
        index = _load_index()
        index_dirty = False
        missing_extractors: set[str] = set()
        for path in sorted(root.rglob("*")):
            if cancelled():
                break
            if not path.is_file():
                continue
            suffix = path.suffix.lower()
            if suffix not in (_TEXT_EXTS | _PDF_EXTS | _PPTX_EXTS | _DOCX_EXTS):
                continue
            if path.name == "content_index.json":
                continue
            try:
                if path.stat().st_size > MAX_FILE_MB * 1024 * 1024:
                    continue
                key = _cache_key(path)
            except OSError:
                continue
            cached = index.get(key)
            if cached is not None:
                text = cached.get("text", "")
                kind = cached.get("kind", "")
            else:
                text, kind = extract_file_text(path)
                if text is None:
                    missing_extractors.add(kind)
                    continue
                index[key] = {"text": text[:MAX_CHARS_PER_FILE], "kind": kind}
                index_dirty = True
            ratio, pos = _match(text or "", terms)
            if ratio >= min_ratio:
                results.append(_result(
                    f"file:{kind}", path.name, str(path), ratio,
                    _snippet(text, pos)))
        if index_dirty:
            # Prune entries for files that no longer exist (bounded size).
            live = {k: v for k, v in index.items()
                    if Path(k.split("|", 1)[0]).exists()}
            _save_index(live)
        for kind in sorted(missing_extractors):
            lib = {"pdf": "pypdf", "pptx": "python-pptx",
                   "docx": "python-docx"}.get(kind, kind)
            skipped.append(f"{kind} files ({lib} not installed)")
    else:
        skipped.append("education files (storage path unavailable)")

    # 5) Consultation notes (permission-gated) ---------------------------
    if include_consultations:
        try:
            from database import consultation_db
            sources += 1
            for row in consultation_db.list_consultations():
                blob = " ".join(str(row.get(k) or "") for k in
                                ("clinical_question", "notes", "title",
                                 "response_text", "summary"))
                ratio, pos = _match(blob, terms)
                if ratio >= min_ratio:
                    results.append(_result(
                        "consultation",
                        str(row.get("title") or row.get("consultation_id")
                            or "Consultation"),
                        f"consultation_id={row.get('consultation_id')}",
                        ratio, _snippet(blob, pos)))
        except Exception:
            logger.exception("content_search: consultation search failed")
            skipped.append("consultations (error)")

    results.sort(key=lambda r: r["score"], reverse=True)
    return {
        "query": query,
        "results": results,
        "sources_searched": sources,
        "skipped_extractors": skipped,
        "elapsed_s": round(time.time() - t0, 2),
    }


__all__ = ["search_education_content", "extract_file_text"]
